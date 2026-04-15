from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

OUTPUT_PATH = "docs/slides/client_server_data_flow.pptx"


def set_title_style(shape):
    tf = shape.text_frame
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)


def style_textbox(text_frame, size=20, color=(240, 240, 240), bold=False):
    for p in text_frame.paragraphs:
        if not p.runs:
            run = p.add_run()
            run.text = p.text
            p.text = ""
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(*color)


def add_bg(slide, rgb=(12, 24, 44)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)


def add_bullets(slide, x, y, w, h, lines, font_size=20):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(236, 240, 245)
    return box


def add_title_only_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_bg(slide, (8, 20, 38))
    slide.shapes.title.text = title
    set_title_style(slide.shapes.title)
    sub = slide.placeholders[1]
    sub.text = subtitle
    style_textbox(sub.text_frame, size=20, color=(210, 220, 235))


def add_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_bg(slide)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7))
    title.text_frame.text = "Repository Client-Server Architecture"
    style_textbox(title.text_frame, size=32, bold=True)

    client = slide.shapes.add_shape(1, Inches(0.6), Inches(1.4), Inches(3.4), Inches(1.8))
    client.fill.solid()
    client.fill.fore_color.rgb = RGBColor(34, 99, 167)
    client.line.color.rgb = RGBColor(130, 180, 235)
    client.text = "React Client\n(client/src)\n\nCalls API helpers\nin api/simulation.js"
    style_textbox(client.text_frame, size=18)

    proxy = slide.shapes.add_shape(1, Inches(4.4), Inches(1.65), Inches(3.2), Inches(1.3))
    proxy.fill.solid()
    proxy.fill.fore_color.rgb = RGBColor(42, 62, 88)
    proxy.line.color.rgb = RGBColor(124, 151, 187)
    proxy.text = "Vite Proxy\n/api -> localhost:3001"
    style_textbox(proxy.text_frame, size=16)

    server = slide.shapes.add_shape(1, Inches(8.0), Inches(1.2), Inches(4.2), Inches(2.0))
    server.fill.solid()
    server.fill.fore_color.rgb = RGBColor(29, 120, 88)
    server.line.color.rgb = RGBColor(136, 218, 188)
    server.text = "Express Server\n(server/src/index.js)\n\nRoutes under /api/*"
    style_textbox(server.text_frame, size=18)

    db = slide.shapes.add_shape(1, Inches(8.0), Inches(3.7), Inches(4.2), Inches(2.0))
    db.fill.solid()
    db.fill.fore_color.rgb = RGBColor(114, 74, 24)
    db.line.color.rgb = RGBColor(228, 181, 107)
    db.text = "Database Module\n(server/src/database/db.js)\n\nIn-memory + JSON persistence\n(server/data/lunar_lander.json)"
    style_textbox(db.text_frame, size=16)

    slide.shapes.add_connector(1, Inches(4.0), Inches(2.3), Inches(4.4), Inches(2.3))
    slide.shapes.add_connector(1, Inches(7.6), Inches(2.3), Inches(8.0), Inches(2.2))
    slide.shapes.add_connector(1, Inches(10.1), Inches(3.2), Inches(10.1), Inches(3.7))

    add_bullets(
        slide,
        0.6,
        5.7,
        11.7,
        1.5,
        [
            "Client UI requests data via axios in client/src/api/simulation.js.",
            "Vite dev server forwards /api calls to Express (port 3001).",
            "Express routes call db queries and physics functions, then return JSON responses.",
        ],
        16,
    )


def add_api_catalog_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_bg(slide, (10, 18, 32))
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7))
    title.text_frame.text = "How The Client Requests Data (API Map)"
    style_textbox(title.text_frame, size=30, bold=True)

    left_lines = [
        "Read endpoints used by UI:",
        "GET /api/constants",
        "GET /api/equations and /api/equations/:category",
        "GET /api/landers and /api/landers/:id",
        "GET /api/simulation/:sessionId",
        "GET /api/simulation/:sessionId/telemetry",
        "GET /api/sessions",
    ]
    right_lines = [
        "State-changing endpoints:",
        "POST /api/simulation/start",
        "POST /api/simulation/:sessionId/step",
        "POST /api/simulation/:sessionId/mode",
        "POST /api/simulation/:sessionId/reset",
        "POST /api/physics/gravity",
        "POST /api/physics/time-to-impact",
    ]

    add_bullets(slide, 0.7, 1.2, 5.8, 5.6, left_lines, 18)
    add_bullets(slide, 6.6, 1.2, 5.8, 5.6, right_lines, 18)

    add_bullets(
        slide,
        0.7,
        6.2,
        11.7,
        0.8,
        ["All endpoints are served by server/src/index.js and consumed in client/src/api/simulation.js."],
        16,
    )


def add_request_flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_bg(slide, (12, 23, 40))
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7))
    title.text_frame.text = "Example Data Flow: Start + Step Simulation"
    style_textbox(title.text_frame, size=30, bold=True)

    flow_lines = [
        "1) Client calls startSimulation(config) -> POST /api/simulation/start.",
        "2) Server loads selected lander via queries.getLanderById(landerId).",
        "3) Server creates session, initializes Blackboard state, logs telemetry.",
        "4) Response returns sessionId, initial state, lander, constants.",
        "5) Client calls stepSimulation(sessionId, thrustPercent, dt).",
        "6) Server advances physics, updates activeSimulations, and logs telemetry.",
        "7) Response returns new state JSON for UI and telemetry panels.",
    ]
    add_bullets(slide, 0.8, 1.2, 11.7, 4.9, flow_lines, 18)

    payload = [
        "Response objects flowing back to client include:",
        "state.altitude, state.velocity, state.fuel, state.gravity, state.timeToImpact,",
        "state.touchdownVelocity, landingResult, mode, and history metadata.",
    ]
    add_bullets(slide, 0.8, 5.7, 11.7, 1.3, payload, 16)


def add_database_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_bg(slide, (16, 20, 28))
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7))
    title.text_frame.text = "How The Database Sends Information Over"
    style_textbox(title.text_frame, size=30, bold=True)

    left = [
        "Reference data synced on startup:",
        "physical_constants",
        "equations",
        "lander_configs",
        "",
        "Runtime session data:",
        "simulation_sessions",
        "telemetry_log",
    ]
    right = [
        "Read path:",
        "queries.getAllConstants/getAllEquations/getAllLanders",
        "queries.getSession/getAllSessions/getTelemetry",
        "",
        "Write path:",
        "queries.createSession",
        "queries.updateSession",
        "queries.logTelemetry",
        "",
        "Persistence:",
        "saveDatabase() writes JSON to server/data/lunar_lander.json",
    ]

    add_bullets(slide, 0.8, 1.2, 5.5, 5.8, left, 18)
    add_bullets(slide, 6.2, 1.2, 6.0, 5.8, right, 16)


def add_sequence_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_bg(slide, (9, 25, 30))
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7))
    title.text_frame.text = "Sequence: Client Request -> API Route -> DB Query -> Client Response"
    style_textbox(title.text_frame, size=26, bold=True)

    sequence = [
        "A. User opens dashboard -> client loads constants, equations, and landers.",
        "B. React component calls getConstants/getEquations/getLanders (axios).",
        "C. Express route handlers call query layer in db.js.",
        "D. Query layer returns JSON rows from in-memory structures.",
        "E. Route sends res.json(...) to client through Vite proxy.",
        "F. React updates telemetry panels and simulation visual state.",
        "G. During simulation, telemetry is logged and can be fetched on demand.",
    ]

    add_bullets(slide, 0.8, 1.3, 11.7, 5.7, sequence, 17)


def add_takeaways_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_bg(slide, (11, 19, 37))
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7))
    title.text_frame.text = "Key Takeaways"
    style_textbox(title.text_frame, size=32, bold=True)

    bullets = [
        "The client never accesses the database directly; all data flows through API routes.",
        "server/src/index.js orchestrates route handling, simulation state, and db query calls.",
        "server/src/database/db.js provides durable JSON-backed storage plus query helpers.",
        "client/src/api/simulation.js centralizes request functions used by React components.",
        "This architecture keeps UI concerns separate from simulation and persistence concerns.",
    ]
    add_bullets(slide, 0.8, 1.3, 11.6, 5.8, bullets, 19)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_only_slide(
        prs,
        "Lunar Lander Repo: Client-Server Data Flow",
        "How frontend requests travel through APIs and return database-backed simulation data",
    )
    add_architecture_slide(prs)
    add_api_catalog_slide(prs)
    add_request_flow_slide(prs)
    add_database_slide(prs)
    add_sequence_slide(prs)
    add_takeaways_slide(prs)

    prs.save(OUTPUT_PATH)
    print(f"Created slide deck: {OUTPUT_PATH}")


if __name__ == "__main__":
    build_presentation()
